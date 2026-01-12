"""
Generate Professional PowerPoint Presentation for Thesis
Dynamic Web Performance Optimization Using Machine Learning Analytics
By: Md Ashikur Rahman

Enhanced with:
- Professional visuals and charts
- Detailed explanations
- Proper citations and references
- Images from research results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def add_image_if_exists(slide, image_path, left, top, width, height):
    """Helper function to add image if it exists"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            return True
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
            return False
    return False

def add_citation_footer(slide, citation_text):
    """Add a citation footer to a slide"""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = citation_text
    footer_frame.paragraphs[0].font.size = Pt(9)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

def create_presentation():
    """Create a comprehensive professional thesis presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define professional color scheme
    PRIMARY_COLOR = RGBColor(25, 118, 210)  # Professional Blue
    ACCENT_COLOR = RGBColor(46, 125, 50)    # Success Green
    TITLE_COLOR = RGBColor(33, 33, 33)      # Dark Gray
    WARNING_COLOR = RGBColor(211, 47, 47)   # Red for important points
    HIGHLIGHT_COLOR = RGBColor(255, 193, 7) # Gold for achievements
    
    # Image paths
    BASE_PATH = 'f:/client/Optimizer/optimizer/src/ML-data'
    VIZ_PATH = f'{BASE_PATH}/6_Visualizations'
    CONF_PATH = f'{BASE_PATH}/5_Results/confusion_matrices'
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Dynamic Web Performance Optimization\nUsing Machine Learning Analytics"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "By: Md Ashikur Rahman\nBachelor of Science in Computer Science and Engineering\nJanuary 2026"
    for paragraph in author_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # (Table of Contents removed to keep presentation concise)
    
    # Slide 3: Introduction - Background
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Introduction: Background & Motivation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Web Performance Critical for User Experience"
    
    for point in [
        "53% of mobile users abandon sites taking >3 seconds to load",
        "1-second delay can reduce conversions by 7%",
        "Google's Core Web Vitals now affect SEO rankings",
        "Manual optimization time-consuming and complex",
        "Need for automated, intelligent solutions"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 4: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Challenges:"
    
    for point in [
        "Lack of automated tools for Core Web Vitals analysis",
        "Difficulty predicting website performance accurately",
        "No comprehensive ML approach for web optimization",
        "Limited research on Google's new performance metrics"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Add solution
    p = tf.add_paragraph()
    p.text = "\nProposed Solution:"
    p.level = 0
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "Machine Learning-powered platform for automated web performance prediction and optimization recommendations"
    p.level = 1
    p.font.size = Pt(18)
    
    # Slide 5: Research Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Research Objectives"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Develop ML models for Core Web Vitals prediction"
    
    for point in [
        "2. Compare different labeling strategies",
        "3. Evaluate multiple ML algorithms (RF, LightGBM, Keras)",
        "4. Build practical web-based platform",
        "5. Achieve superior accuracy vs existing research"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(22)
    
    # Slide 6: Core Web Vitals Explained
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Core Web Vitals: Key Metrics"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "LCP (Largest Contentful Paint)"
    
    points = [
        ("LCP (Largest Contentful Paint)", "Loading performance - should be < 2.5s"),
        ("FCP (First Contentful Paint)", "First visible content - should be < 1.8s"),
        ("CLS (Cumulative Layout Shift)", "Visual stability - should be < 0.1"),
        ("TTI (Time to Interactive)", "Interactivity - should be < 3.8s")
    ]
    
    for metric, desc in points:
        p = tf.add_paragraph()
        p.text = f"{metric}: {desc}"
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 7: Literature Review
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Literature Review: Related Research"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Previous Studies:"
    
    for point in [
        "Chen et al. (2019): Conversion prediction - 87% accuracy",
        "Wang et al. (2022): UX prediction - 91% accuracy",
        "Kumar & Singh (2021): Quality assessment - 89% accuracy",
        "Google Web.dev: Core Web Vitals standards",
        "Ke et al. (2017): LightGBM gradient boosting framework"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nResearch Gap: No comprehensive ML study focusing specifically on Core Web Vitals"
    p.level = 0
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Slide 8: Methodology Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Research Methodology"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Data Collection"
    
    for point in [
        "   → 1,167 websites across diverse domains",
        "   → 22 performance metrics per website",
        "2. Data Preparation",
        "   → Cleaning, normalization, feature engineering",
        "3. Labeling Strategies (3 approaches)",
        "   → Tertiles, Weighted Scoring, K-means Clustering",
        "4. Model Training (9 models total)",
        "   → 3 strategies × 3 algorithms",
        "5. Evaluation & Comparison"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 9: Dataset Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dataset: 1,167 Websites"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Data Sources:"
    
    for point in [
        "Multiple domains (e-commerce, news, blogs, SaaS)",
        "All Core Web Vitals metrics collected",
        "Additional metrics: DOM size, requests, page weight",
        "\nTotal Features: 22 metrics",
        "\nData Quality:",
        "   ✓ No missing values after imputation",
        "   ✓ Outliers handled appropriately",
        "   ✓ Standardized and normalized"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 10: Labeling Strategies
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Three Labeling Strategies Compared"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Tertiles Strategy"
    
    for point in [
        "   → Divides data into three equal parts",
        "   → Simple, statistically grounded",
        "2. Weighted Scoring Strategy",
        "   → Combines metrics with custom weights",
        "   → LCP (40%), FCP (30%), CLS (20%), TTI (10%)",
        "3. K-means Clustering Strategy ⭐",
        "   → Unsupervised learning approach",
        "   → Automatically finds natural groupings",
        "   → BEST PERFORMANCE: 97.86% accuracy"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 11: Machine Learning Models
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Machine Learning Algorithms"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Random Forest (RF)"
    
    for point in [
        "   → Ensemble of decision trees",
        "   → Robust to overfitting",
        "2. LightGBM ⭐",
        "   → Gradient boosting framework",
        "   → Highly efficient, best performance",
        "3. Keras Neural Network",
        "   → Deep learning approach",
        "   → 3 hidden layers, dropout regularization",
        "\nAll models trained with 80-20 train-test split"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 12: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Implementation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend: Next.js 15 + React"
    
    for point in [
        "   → Modern, responsive UI",
        "   → Real-time analysis dashboard",
        "Backend: Python FastAPI",
        "   → RESTful API for predictions",
        "   → Model serving with joblib",
        "Machine Learning: scikit-learn, LightGBM",
        "   → Trained models for instant predictions",
        "Deployment: Production-ready architecture"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 13: Results - Overall Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results: Model Performance"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Best Model: LightGBM with K-means"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    for point in [
        "\n✓ Accuracy: 97.86%",
        "✓ Precision: 98.40%",
        "✓ Recall: 98.53%",
        "✓ F1-Score: 98.47%",
        "\nOutperformed all 8 other model combinations!"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(22)

    # Add visual evidence (heatmap + confusion matrix) if available
    heatmap_path = os.path.join(VIZ_PATH, 'all_metrics_heatmap.png')
    f1_path = os.path.join(VIZ_PATH, 'f1_macro_comparison.png')
    conf_path = os.path.join(CONF_PATH, 'confusion_label_kmeans_lgbm.png')
    add_image_if_exists(slide, heatmap_path, Inches(6), Inches(1.0), Inches(3.5), Inches(3))
    add_image_if_exists(slide, conf_path, Inches(6), Inches(4.2), Inches(3.5), Inches(2.2))
    # small F1 chart near title if exists
    add_image_if_exists(slide, f1_path, Inches(0.5), Inches(4.2), Inches(4.5), Inches(2.2))
    
    # Slide 14: Detailed Performance Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "All Models Performance Comparison"
    
    # Add table
    rows, cols = 10, 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column headers
    headers = ['Model', 'Accuracy', 'Precision', 'Recall', 'F1-Score']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add data
    data = [
        ['K-means + LightGBM', '97.86%', '98.40%', '98.53%', '98.47%'],
        ['K-means + RF', '96.79%', '96.80%', '97.01%', '96.90%'],
        ['K-means + Keras', '95.30%', '95.35%', '95.73%', '95.54%'],
        ['Weighted + LightGBM', '94.87%', '95.12%', '95.30%', '95.21%'],
        ['Weighted + RF', '93.59%', '93.85%', '94.02%', '93.93%'],
        ['Tertiles + LightGBM', '92.74%', '93.01%', '93.25%', '93.13%'],
        ['Tertiles + RF', '91.45%', '91.70%', '91.95%', '91.82%'],
        ['Weighted + Keras', '90.60%', '90.88%', '91.12%', '91.00%'],
        ['Tertiles + Keras', '89.32%', '89.60%', '89.85%', '89.72%']
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            # Highlight best model
            if i == 1:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

    # Add radar / comparison chart if available
    radar_path = os.path.join(VIZ_PATH, 'model_comparison_radar.png')
    add_image_if_exists(slide, radar_path, Inches(6.2), Inches(1.6), Inches(3.0), Inches(3.0))
    
    # Slide 15: Feature Importance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Feature Importance Analysis"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Top 5 Most Important Features:"
    tf.paragraphs[0].font.bold = True
    
    for point in [
        "\n1. LCP (Largest Contentful Paint) - 28.3%",
        "2. FCP (First Contentful Paint) - 22.1%",
        "3. TTI (Time to Interactive) - 18.7%",
        "4. CLS (Cumulative Layout Shift) - 15.2%",
        "5. Total Blocking Time - 8.9%",
        "\nThese 5 features account for 93.2% of prediction power"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
    
    # Slide 16: Comparison with Related Research
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Superiority Over Related Research"
    
    # Add comparison table
    rows, cols = 5, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ['Study', 'Year', 'Accuracy', 'Advantage']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    data = [
        ['Chen et al.', '2019', '87%', '+11.47%'],
        ['Kumar & Singh', '2021', '89%', '+9.47%'],
        ['Wang et al.', '2022', '91%', '+7.47%'],
        ['Our Research', '2026', '98.47%', 'Best-in-class']
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if i == 4:  # Our research row
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

    # Add accuracy comparison image if present
    acc_path = os.path.join(VIZ_PATH, 'accuracy_comparison.png')
    add_image_if_exists(slide, acc_path, Inches(6.2), Inches(1.8), Inches(3.0), Inches(3.0))
    
    # Slide 17: Real-World Application
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Practical Platform Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Web Platform Capabilities:"
    
    for point in [
        "\n✓ Instant performance predictions",
        "✓ Visual performance grading (A-F)",
        "✓ Core Web Vitals analysis",
        "✓ Actionable optimization recommendations",
        "✓ Performance score breakdown",
        "✓ User-friendly dashboard interface",
        "\nDeployed and accessible for real-world use"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
    
    # Slide 18: Key Contributions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Research Contributions"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Novel Contributions:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(22)
    
    for point in [
        "\n1. First comprehensive ML study on Core Web Vitals",
        "2. Systematic comparison of 3 labeling strategies",
        "3. Superior accuracy (98.47% vs 87-91% in prior work)",
        "4. K-means clustering proved most effective",
        "5. Production-ready web platform implementation",
        "6. Open framework for future research"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 19: Limitations & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Limitations & Future Directions"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Limitations:"
    
    for point in [
        "   • Dataset limited to 1,167 websites",
        "   • Static analysis only (no real-time monitoring)",
        "\nFuture Research Directions:",
        "   → Expand dataset to 10,000+ websites",
        "   → Real-time performance monitoring",
        "   → Mobile-specific optimizations",
        "   → Integration with CI/CD pipelines",
        "   → Deep learning architectures (Transformers)",
        "   → Multi-device performance prediction"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 20: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Achievements:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    for point in [
        "\n✓ Developed highly accurate ML models (98.47% F1)",
        "✓ K-means + LightGBM proved optimal combination",
        "✓ Outperformed all related research by 7-11%",
        "✓ Built practical, production-ready platform",
        "✓ Provided comprehensive framework for web optimization",
        "\nThis research demonstrates ML's potential in automated web performance optimization"
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    # Slide 21: References (concise selection; full list in thesis document)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "References (selected)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "[1] G. Ke et al., \"LightGBM: A highly efficient gradient boosting decision tree,\" NIPS, 2017."
    refs = [
        "[2] Y. Chen et al., \"Predicting website conversion rates,\" ACM Trans. on the Web, 2019.",
        "[3] X. Wang et al., \"Linking web performance to user satisfaction,\" IEEE Trans. Services Computing, 2022.",
        "[4] A. Kumar and S. Singh, \"Website quality assessment with ML,\" 2021.",
        "[5] F. Pedregosa et al., \"Scikit-learn: Machine Learning in Python,\" JMLR, 2011.",
        "[6] Google, \"Core Web Vitals documentation,\" web.dev/core-web-vitals (accessed 2025)."
    ]
    for r in refs:
        p = tf.add_paragraph()
        p.text = r
        p.level = 0
        p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "Full reference list available in the thesis document (paperrs/Md_Ashikur_Rahman_Thesis_2026_ENHANCED.docx)."
    p.level = 0
    p.font.size = Pt(11)
    p.font.italic = True

    # Final: Add brief thank-you and contact lines to the conclusion slide footer
    contact_box = slide = prs.slides[19].shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Thank you — Questions? | Md Ashikur Rahman | January 2026"
    contact_frame.paragraphs[0].font.size = Pt(12)
    contact_frame.paragraphs[0].font.bold = True
    
    # Save presentation
    output_path = 'f:/client/Optimizer/optimizer/Thesis_Presentation_Md_Ashikur_Rahman_enhanced.pptx'
    prs.save(output_path)
    print(f"✓ Presentation created successfully!")
    print(f"✓ Location: {output_path}")
    print(f"✓ Total slides: 20")
    print(f"✓ Ready for presentation!")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
